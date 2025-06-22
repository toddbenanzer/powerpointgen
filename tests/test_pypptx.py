import unittest
import os
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER

# Assuming pypptx is installed or PYTHONPATH is set up correctly
# For local testing, you might need to adjust sys.path
import sys
# Add project root to sys.path to allow importing pypptx
# This assumes the tests are run from the project root directory or similar context
# For a robust solution, packaging and installation is preferred.
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from pypptx.presentation import PyPPT
from pypptx.constants import (
    DEFAULT_SUBTITLE_FONT_NAME,
    DEFAULT_SUBTITLE_FONT_SIZE_PT,
    DEFAULT_FOOTER_FONT_NAME,
    DEFAULT_FOOTER_FONT_SIZE_PT
)

class TestPyPPTXFormatting(unittest.TestCase):

    def setUp(self):
        """Set up a new presentation for each test."""
        self.ppt = PyPPT()
        # Use a layout that is likely to have a subtitle (e.g., Title Slide - layout 0)
        # For footers, they are often controlled by the slide master and can appear on many layouts.
        self.slide = self.ppt.add_slide(layout_ref=0) # Title Slide layout

    def find_placeholder(self, slide_obj, placeholder_type):
        for shape in slide_obj.pptx_slide.placeholders:
            if shape.placeholder_format.type == placeholder_type:
                return shape
        return None

    def test_set_subtitle_applies_default_formatting(self):
        """Test that set_subtitle applies text and default formatting."""
        subtitle_text = "This is a Test Subtitle"

        # The `set_subtitle` method in PySlide looks for PP_PLACEHOLDER.SUBTITLE.
        # We need to ensure our test slide (from layout 0) actually has one.
        # A default 'Title Slide' (layout 0) usually has a title placeholder (idx 0)
        # and a subtitle placeholder (idx 1, often of type PP_PLACEHOLDER.BODY or PP_PLACEHOLDER.SUBTITLE).
        # If it's not specifically PP_PLACEHOLDER.SUBTITLE, set_subtitle will fail.

        target_placeholder = self.find_placeholder(self.slide, PP_PLACEHOLDER.SUBTITLE)

        if not target_placeholder:
            # Fallback: Check if placeholder at index 1 (common for subtitle on title slide) exists
            # and if it's a type that could accept text (e.g. BODY, OBJECT, or even a generic one)
            # This is to make the test more robust if the default layout 0 doesn't use strict SUBTITLE type.
            # However, set_subtitle is strict. So, this test *must* find PP_PLACEHOLDER.SUBTITLE.
            self.skipTest(f"Slide layout 0 does not have a placeholder of type PP_PLACEHOLDER.SUBTITLE. " +
                          "The set_subtitle method requires this specific type.")

        try:
            self.slide.set_subtitle(subtitle_text)
        except AttributeError as e:
            self.fail(f"set_subtitle raised an AttributeError: {e}. Layout 0 might lack a true PP_PLACEHOLDER.SUBTITLE.")

        # Re-fetch the placeholder that set_subtitle would have acted upon.
        updated_subtitle_shape = self.find_placeholder(self.slide, PP_PLACEHOLDER.SUBTITLE)

        self.assertIsNotNone(updated_subtitle_shape, "Subtitle placeholder of type SUBTITLE should exist after set_subtitle.")
        self.assertTrue(updated_subtitle_shape.has_text_frame, "Subtitle placeholder should have a text frame.")
        # Ensure text was actually set in the correct placeholder
        self.assertEqual(updated_subtitle_shape.text_frame.text, subtitle_text, "Subtitle text not set correctly in the SUBTITLE placeholder.")

        self.assertTrue(len(updated_subtitle_shape.text_frame.paragraphs) > 0, "Subtitle text frame should have paragraphs.")
        font = updated_subtitle_shape.text_frame.paragraphs[0].font
        self.assertEqual(font.name, DEFAULT_SUBTITLE_FONT_NAME, "Subtitle font name not set to default.")
        self.assertEqual(font.size, Pt(DEFAULT_SUBTITLE_FONT_SIZE_PT), "Subtitle font size not set to default.")

    def test_set_footer_applies_default_formatting(self):
        """Test that set_footer_text applies text and default formatting."""
        footer_text = "This is a Test Footer"

        target_footer_placeholder = self.find_placeholder(self.slide, PP_PLACEHOLDER.FOOTER)

        if not target_footer_placeholder:
            self.skipTest(f"Slide layout 0 does not have an explicit PP_PLACEHOLDER.FOOTER. " +
                          "This test requires a layout/master with an active footer placeholder.")

        try:
            self.slide.set_footer_text(footer_text)
        except AttributeError as e:
            self.fail(f"set_footer_text raised an AttributeError: {e}. The layout may lack a footer placeholder.")

        updated_footer_shape = self.find_placeholder(self.slide, PP_PLACEHOLDER.FOOTER)

        self.assertIsNotNone(updated_footer_shape, "Footer placeholder of type FOOTER should exist after set_footer_text.")
        self.assertTrue(updated_footer_shape.has_text_frame, "Footer placeholder should have a text frame.")
        self.assertEqual(updated_footer_shape.text_frame.text, footer_text, "Footer text not set correctly in the FOOTER placeholder.")

        self.assertTrue(len(updated_footer_shape.text_frame.paragraphs) > 0, "Footer text frame should have paragraphs.")
        font = updated_footer_shape.text_frame.paragraphs[0].font
        self.assertEqual(font.name, DEFAULT_FOOTER_FONT_NAME, "Footer font name not set to default.")
        self.assertEqual(font.size, Pt(DEFAULT_FOOTER_FONT_SIZE_PT), "Footer font size not set to default.")

if __name__ == '__main__':
    unittest.main()
