from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd

# Re-export enums for easier access by users
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER # Consolidated
from pptx.enum.chart import XL_CHART_TYPE

__all__ = ['PyPPT', 'PySlide', 'MSO_SHAPE', 'XL_CHART_TYPE', 'PP_PLACEHOLDER']

# Default values for method parameters
DEFAULT_LAYOUT_REF = 5
DEFAULT_TABLE_HEADER_BOLD = True
DEFAULT_TABLE_INCLUDE_INDEX = False

class PySlide:
    def __init__(self, pptx_slide):
        """Initializes the PySlide.

        Args:
            pptx_slide (pptx.slide.Slide): The python-pptx Slide object.
        """
        self.pptx_slide = pptx_slide

    def set_title(self, text):
        """Sets the title of this slide.

        Args:
            text (str): The text to set as the title.

        Raises:
            AttributeError: If the slide does not have a title placeholder.
        """
        if self.pptx_slide.shapes.title:
            self.pptx_slide.shapes.title.text = text
        else:
            raise AttributeError("This slide does not have a title placeholder.")

    def set_subtitle(self, text):
        """Sets the subtitle of this slide.

        Args:
            text (str): The text to set as the subtitle.

        Raises:
            AttributeError: If the slide does not have a suitable subtitle placeholder.
        """
        subtitle_shape = None
        for shape in self.pptx_slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
                subtitle_shape = shape
                break

        if subtitle_shape:
            subtitle_shape.text = text
        else:
            raise AttributeError("This slide does not have a clear subtitle placeholder.")

    def set_footer_text(self, text):
        """Sets the footer text on this slide.

        Args:
            text (str): The text to set in the footer.

        Raises:
            AttributeError: If the slide does not have a footer placeholder.
        """
        footer_shape = None
        for shape in self.pptx_slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                footer_shape = shape
                break

        if footer_shape:
            footer_shape.text = text
        else:
            raise AttributeError("This slide does not have a footer placeholder.")

    def add_text_box(self, text, left, top, width, height):
        """Adds a text box with plain text to this slide.
        Positions and dimensions are in Inches.

        Args:
            text (str): The plain text to add to the text box.
            left (float): The left position of the text box (Inches).
            top (float): The top position of the text box (Inches).
            width (float): The width of the text box (Inches).
            height (float): The height of the text box (Inches).
        """
        shape = self.pptx_slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        shape.text_frame.text = text
        return shape

    def add_bullet_point_box(self, items, left, top, width, height):
        """Adds a text box with a list of items as bullet points to this slide.
        Positions and dimensions are in Inches.

        Args:
            items (list[str]): A list of strings, each will be a bullet point.
            left (float): The left position (in Inches).
            top (float): The top position (in Inches).
            width (float): The width (in Inches).
            height (float): The height (in Inches).
        """
        left_emu = Inches(left)
        top_emu = Inches(top)
        width_emu = Inches(width)
        height_emu = Inches(height)

        shape = self.pptx_slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
        tf = shape.text_frame
        tf.clear()

        for item_text in items:
            p = tf.add_paragraph()
            p.text = item_text
            p.level = 0

        return shape

    def add_table_from_dataframe(self, dataframe, left, top, width, height,
                                 column_labels=None, number_formats=None,
                                 include_index=DEFAULT_TABLE_INCLUDE_INDEX, index_label=None,
                                 font_name=None, font_size=None,
                                 column_widths=None,
                                 row_heights=None,
                                 header_bold=DEFAULT_TABLE_HEADER_BOLD,
                                 header_font_color_rgb=None,
                                 header_fill_color_rgb=None
                                 ):
        """Adds a table to the slide populated from a Pandas DataFrame with styling.

        Args:
            dataframe (pd.DataFrame): The Pandas DataFrame to display.
            left (float): Left position of the table (Inches).
            top (float): Top position of the table (Inches).
            width (float): Width of the table (Inches).
            height (float): Height of the table (Inches).
            column_labels (dict, optional): Maps DataFrame column names to display names.
            number_formats (dict, optional): Maps DataFrame column names to format strings.
            include_index (bool): True to include DataFrame index as first column. Defaults to DEFAULT_TABLE_INCLUDE_INDEX.
            index_label (str, optional): Header for the index column.
            font_name (str, optional): Font name for table text (e.g., "Arial").
            font_size (int, optional): Font size for table text (in Points, e.g., 10).
            column_widths (list|dict, optional): List or dict of column widths in Inches.
                                               If list, applied by index. If dict, by col_idx.
            row_heights (list|dict, optional): List or dict of row heights in Inches.
                                             If list, applied by index. If dict, by row_idx.
            header_bold (bool): True to make header text bold. Defaults to DEFAULT_TABLE_HEADER_BOLD.
            header_font_color_rgb (tuple, optional): RGB tuple for header font color (e.g., (255,255,255)).
            header_fill_color_rgb (tuple, optional): RGB tuple for header cell fill color (e.g., (0,0,0)).
        Returns:
            pptx.shapes.graphfrm.GraphicFrame: The table shape object.
        """
        if not isinstance(dataframe, pd.DataFrame):
            raise ValueError("Input 'dataframe' must be a pandas DataFrame.")

        rows = len(dataframe) + 1  # +1 for header row
        cols = len(dataframe.columns)
        if include_index:
            cols += 1

        # Create the table shape
        table_shape = self.pptx_slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        table = table_shape.table

        # --- Populate Header Row ---
        current_col_idx = 0
        if include_index:
            header_text = index_label if index_label is not None else (dataframe.index.name if dataframe.index.name is not None else "Index")
            table.cell(0, current_col_idx).text = str(header_text)
            current_col_idx += 1

        for df_col_name in dataframe.columns:
            display_name = str(df_col_name) # Default to original column name
            if column_labels and df_col_name in column_labels:
                display_name = str(column_labels[df_col_name])
            table.cell(0, current_col_idx).text = display_name
            current_col_idx += 1

        # --- Populate Data Rows ---
        for i, df_row_tuple in enumerate(dataframe.itertuples(index=include_index, name=None)):
            # df_row_tuple contains actual data values. If include_index=True, first element is index.

            current_cell_idx_in_row = 0
            data_tuple_offset = 0 # if include_index is False, df_row_tuple starts with first data col

            if include_index:
                index_val = df_row_tuple[0]
                # Apply formatting to index if 'index_label' (or a convention) is in number_formats
                # For now, index is added as string without specific formatting via number_formats
                table.cell(i + 1, current_cell_idx_in_row).text = str(index_val)
                current_cell_idx_in_row += 1
                data_tuple_offset = 1 # Data starts from 2nd element of df_row_tuple

            for col_idx, df_col_name in enumerate(dataframe.columns):
                cell_value = df_row_tuple[col_idx + data_tuple_offset]

                formatted_value = str(cell_value) # Default to string representation
                if pd.isna(cell_value): # Handle NaN/None before formatting
                    formatted_value = "" # Or some other placeholder like "N/A"
                elif number_formats and df_col_name in number_formats:
                    try:
                        fmt_spec = number_formats[df_col_name]
                        # Ensure value is appropriate for format spec (e.g. numeric for 'f', 'd')
                        if isinstance(cell_value, (int, float)):
                             formatted_value = f"{cell_value:{fmt_spec}}"
                        # else: formatted_value remains str(cell_value) if format spec is for numbers but type isn't
                    except ValueError:
                        # print(f"Warning: Could not apply format '{fmt_spec}' to value '{cell_value}' in col '{df_col_name}'. Using default string.")
                        pass # Keep default string if formatting fails

                table.cell(i + 1, current_cell_idx_in_row).text = formatted_value
                current_cell_idx_in_row += 1

        # --- Apply Table-wide Font Styling ---
        if font_name or font_size:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if font_name:
                            paragraph.font.name = font_name
                        if font_size:
                            paragraph.font.size = Pt(font_size)

        # --- Apply Column Widths ---
        if column_widths:
            if isinstance(column_widths, list):
                for i, cw_val in enumerate(column_widths):
                    if i < len(table.columns):
                        table.columns[i].width = Inches(cw_val)
            elif isinstance(column_widths, dict):
                for col_idx, cw_val in column_widths.items():
                    if col_idx < len(table.columns):
                        table.columns[col_idx].width = Inches(cw_val)

        # --- Apply Row Heights ---
        if row_heights:
            if isinstance(row_heights, list):
                for i, rh_val in enumerate(row_heights):
                    if i < len(table.rows):
                        table.rows[i].height = Inches(rh_val)
            elif isinstance(row_heights, dict):
                for row_idx, rh_val in row_heights.items():
                    if row_idx < len(table.rows):
                        table.rows[row_idx].height = Inches(rh_val)

        # --- Style Header Row ---
        for col_idx in range(len(table.columns)):
            cell = table.cell(0, col_idx) # Header row is index 0

            for paragraph in cell.text_frame.paragraphs:
                if header_bold:
                    paragraph.font.bold = True
                if header_font_color_rgb:
                    paragraph.font.color.rgb = RGBColor(*header_font_color_rgb)

            if header_fill_color_rgb:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*header_fill_color_rgb)

        return table_shape

    def add_shape(self, shape_type, left, top, width, height, shape_name=None):
        """Adds a predefined shape to the slide.

        Args:
            shape_type (MSO_SHAPE): The type of shape to add (e.g., MSO_SHAPE.RECTANGLE).
            left (float): The left position of the shape (in Inches).
            top (float): The top position of the shape (in Inches).
            width (float): The width of the shape (in Inches).
            height (float): The height of the shape (in Inches).
            shape_name (str, optional): An optional name for the shape.

        Returns:
            pptx.shape.Shape: The newly added shape object.
        """
        # Ensure Inches is available (should be imported at module level)
        new_shape = self.pptx_slide.shapes.add_shape(
            shape_type,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        if shape_name:
            new_shape.name = shape_name

        return new_shape

    def _get_shape(self, shape_ref):
        """Internal helper to retrieve a shape object.

        Args:
            shape_ref: Can be a shape object itself, the name of a shape (str),
                       or the index of a shape in the slide's shape collection (int).

        Returns:
            pptx.shape.Shape: The found shape object.

        Raises:
            TypeError: If shape_ref is not a Shape, str, or int.
            ValueError: If shape name or index is not found.
        """
        if hasattr(shape_ref, 'shape_type'): # Check if it's already a Shape object (duck typing)
            return shape_ref
        elif isinstance(shape_ref, str): # Find by name
            for shape_in_slide in self.pptx_slide.shapes:
                if shape_in_slide.name == shape_ref:
                    return shape_in_slide
            raise ValueError(f"Shape with name '{shape_ref}' not found on this slide.")
        elif isinstance(shape_ref, int): # Find by index
            try:
                return self.pptx_slide.shapes[shape_ref]
            except IndexError:
                raise ValueError(f"Shape at index {shape_ref} not found on this slide (max index: {len(self.pptx_slide.shapes)-1}).")
        else:
            raise TypeError(f"shape_ref must be a shape object, name (str), or index (int), not {type(shape_ref)}.")

    def set_shape_fill_color(self, shape_ref, r, g, b):
        """Sets the solid fill color of a specified shape.

        Args:
            shape_ref: The shape object, its name (str), or index (int).
            r (int): Red component of RGB color (0-255).
            g (int): Green component of RGB color (0-255).
            b (int): Blue component of RGB color (0-255).
        """
        shape = self._get_shape(shape_ref)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)

    def set_shape_line_color(self, shape_ref, r, g, b):
        """Sets the line color of a specified shape.

        Args:
            shape_ref: The shape object, its name (str), or index (int).
            r (int): Red component of RGB color (0-255).
            g (int): Green component of RGB color (0-255).
            b (int): Blue component of RGB color (0-255).
        """
        shape = self._get_shape(shape_ref)
        shape.line.color.rgb = RGBColor(r, g, b)
        # To make line visible if it was set to 'no line', you might also need:
        # shape.line.visible = True # Or shape.line.width > 0
        # For now, this just sets color. User must ensure line is visible.

    def set_shape_line_weight(self, shape_ref, weight_pt):
        """Sets the line weight (thickness) of a specified shape.

        Args:
            shape_ref: The shape object, its name (str), or index (int).
            weight_pt (float or int): Line weight in points.
        """
        shape = self._get_shape(shape_ref)
        shape.line.width = Pt(weight_pt)
        # Setting width usually makes the line visible if it had 'no line' previously.

    def add_chart(self, chart_type, chart_data_dict, left, top, width, height, chart_title=None):
        """Adds a chart to the slide.

        Args:
            chart_type (XL_CHART_TYPE): Type of chart (e.g., XL_CHART_TYPE.LINE).
            chart_data_dict (dict): Data for the chart. Expected structure:
                {
                    'categories': ['Cat1', 'Cat2', ...],
                    'series': [
                        {'name': 'Series1 Name', 'values': [val1, val2, ...]},
                        {'name': 'Series2 Name', 'values': [valA, valB, ...]}
                    ]
                }
            left (float): Left position of the chart (Inches).
            top (float): Top position of the chart (Inches).
            width (float): Width of the chart (Inches).
            height (float): Height of the chart (Inches).
            chart_title (str, optional): Title for the chart.

        Returns:
            pptx.shapes.graphfrm.GraphicFrame: The graphic frame containing the chart.

        Raises:
            ValueError: If chart_data_dict structure is invalid.
        """
        # Validate chart_data_dict structure
        if not isinstance(chart_data_dict, dict):
            raise ValueError("chart_data_dict must be a dictionary.")
        if 'categories' not in chart_data_dict or not isinstance(chart_data_dict['categories'], list):
            raise ValueError("chart_data_dict must contain a 'categories' list.")
        if 'series' not in chart_data_dict or not isinstance(chart_data_dict['series'], list):
            raise ValueError("chart_data_dict must contain a 'series' list.")
        if not chart_data_dict['series']: # Check if series list is empty
             raise ValueError("chart_data_dict['series'] list cannot be empty.")
        for s in chart_data_dict['series']:
            if not isinstance(s, dict) or 'name' not in s or 'values' not in s:
                raise ValueError("Each item in 'series' must be a dict with 'name' and 'values'.")
            if not isinstance(s['values'], list):
                raise ValueError("Each series 'values' must be a list.")
            if len(s['values']) != len(chart_data_dict['categories']):
                raise ValueError(
                    f"Series '{s['name']}' has {len(s['values'])} values, "
                    f"but there are {len(chart_data_dict['categories'])} categories."
                )

        chart_data = CategoryChartData()
        chart_data.categories = chart_data_dict['categories']

        for series_item in chart_data_dict['series']:
            chart_data.add_series(series_item['name'], series_item['values'])

        graphic_frame = self.pptx_slide.shapes.add_chart(
            chart_type,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        )

        chart = graphic_frame.chart

        if chart_title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_title
        else:
            chart.has_title = False

        if len(chart_data_dict['series']) > 1:
            chart.has_legend = True
        else:
            chart.has_legend = False

        return graphic_frame

class PyPPT:
    def __init__(self, pptx_path=None):
        if pptx_path:
            self.presentation = Presentation(pptx_path)
        else:
            self.presentation = Presentation()
        # Store the path if provided, though it's less relevant if creating new
        self.pptx_path = pptx_path

    def add_slide(self, layout_ref=DEFAULT_LAYOUT_REF):
        """Adds a new slide to the presentation and returns its PySlide wrapper.

        Args:
            layout_ref (int or str): The index of the slide layout to use (int)
                                     or the name of the slide layout (str).
                                     Defaults to DEFAULT_LAYOUT_REF.

        Returns:
            PySlide: A wrapper for the newly added slide.

        Raises:
            ValueError: If layout_ref is a string and no layout with that name is found.
            IndexError: If layout_ref is an integer and is out of range.
            TypeError: If layout_ref is not an int or str.
        """
        slide_layout = None
        if isinstance(layout_ref, str):
            layout_name_to_find = layout_ref
            found_layout = None
            for layout in self.presentation.slide_layouts:
                if hasattr(layout, 'name') and layout.name == layout_name_to_find: # Check if layout has a name attribute
                    found_layout = layout
                    break
            if found_layout:
                slide_layout = found_layout
            else:
                available_layout_names = [l.name for l in self.presentation.slide_layouts if hasattr(l, 'name')]
                raise ValueError(
                    f"Layout with name '{layout_name_to_find}' not found. "
                    f"Available layout names are: {available_layout_names}"
                )
        elif isinstance(layout_ref, int):
            try:
                slide_layout = self.presentation.slide_layouts[layout_ref]
            except IndexError:
                raise IndexError(
                    f"Layout index {layout_ref} is out of range. "
                    f"Available layouts: {len(self.presentation.slide_layouts)}."
                )
        else:
            raise TypeError(f"layout_ref must be an integer index or a string name, not {type(layout_ref)}.")

        new_pptx_slide = self.presentation.slides.add_slide(slide_layout)
        return PySlide(new_pptx_slide)

    def get_slide(self, slide_index):
        """Gets the slide at the specified index as a PySlide instance.

        Args:
            slide_index (int): The index of the slide to retrieve.

        Returns:
            PySlide: A wrapper for the specified slide.

        Raises:
            IndexError: If slide_index is out of range.
        """
        if slide_index < 0 or slide_index >= len(self.presentation.slides):
            raise IndexError(f"Slide index {slide_index} is out of range.")
        pptx_slide = self.presentation.slides[slide_index]
        return PySlide(pptx_slide)

    @property
    def slides(self):
        """Returns a list of PySlide instances for all slides in the presentation."""
        return [PySlide(s) for s in self.presentation.slides]

    def delete_slide(self, slide_index):
        """Deletes a slide from the presentation by its index.

        Args:
            slide_index (int): The index of the slide to delete.

        Raises:
            IndexError: If slide_index is out of range.

        Note: This method manipulates internal structures of python-pptx
              and should be used with caution.
        """
        # Use len(self.presentation.slides) for public API consistency if available
        # but _sldIdLst is the direct list being manipulated.
        num_slides_in_list = len(self.presentation.slides._sldIdLst)
        if not 0 <= slide_index < num_slides_in_list:
            raise IndexError(f"Slide index {slide_index} is out of range. "
                             f"Presentation has {num_slides_in_list} slides (indices 0 to {num_slides_in_list-1}).")

        prs = self.presentation

        slide_id_entry = prs.slides._sldIdLst[slide_index]
        rId = slide_id_entry.rId

        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[slide_index]

    def move_slide(self, current_index, new_index):
        """Moves a slide from its current position to a new position.

        Args:
            current_index (int): The current index of the slide to move.
            new_index (int): The target index where the slide should be moved.
                             If new_index is greater than or equal to the number
                             of slides (after removal of the slide at current_index),
                             the slide is moved to the end.
                             If new_index is less than 0, it's treated as 0.

        Raises:
            IndexError: If current_index is out of range.

        Note: This method manipulates internal structures of python-pptx.
        """
        slides_list = self.presentation.slides._sldIdLst
        num_slides = len(slides_list)

        if not 0 <= current_index < num_slides:
            raise IndexError(f"current_index {current_index} is out of range. "
                             f"Presentation has {num_slides} slides (indices 0 to {num_slides-1}).")

        # Pop the slide ID entry from its current position
        slide_id_entry_to_move = slides_list.pop(current_index)

        # After pop, the list is one shorter.
        # Adjust new_index to be within valid bounds for insertion into the modified list.
        num_slides_after_pop = len(slides_list) # This is num_slides - 1

        if new_index < 0:
            new_index = 0
        elif new_index > num_slides_after_pop: # If new_index is beyond the new end of list
            new_index = num_slides_after_pop   # Clamp to the end (append)

        slides_list.insert(new_index, slide_id_entry_to_move)

    def duplicate_slide(self, slide_index_to_duplicate):
        """Duplicates a slide by creating a new slide with the same layout
        and attempting to copy basic content and shapes.
        The new slide is added at the end of the presentation.

        IMPORTANT LIMITATIONS (Basic Implementation):
        - Does NOT perform a perfect, deep copy of the slide.
        - Copies text from the main title placeholder if present on both slides.
        - For other common placeholders (body, content, text box type), text content
          is copied into NEW text boxes on the duplicated slide at the same position
          and size as the original placeholder. It does not attempt to map to
          existing placeholders on the new slide by index.
        - Attempts to replicate basic auto-shapes (e.g., RECTANGLE, OVAL, LINE)
          and their text content, position, and size.
        - Does NOT copy:
            - Complex shape formatting (most fill, line, effects will be default).
            - Tables.
            - Charts.
            - Images.
            - Grouped shapes.
            - Animations, transitions, or slide master modifications.

        Args:
            slide_index_to_duplicate (int): The index of the slide to duplicate.

        Returns:
            PySlide: A PySlide wrapper for the newly created (duplicated) slide.

        Raises:
            IndexError: If slide_index_to_duplicate is out of range.
        """
        num_slides = len(self.presentation.slides)
        if not 0 <= slide_index_to_duplicate < num_slides:
            raise IndexError(f"slide_index_to_duplicate {slide_index_to_duplicate} is out of range. "
                             f"Presentation has {num_slides} slides (indices 0 to {num_slides-1}).")

        source_slide_pptx = self.presentation.slides[slide_index_to_duplicate]
        source_layout = source_slide_pptx.slide_layout

        new_slide_pptx = self.presentation.slides.add_slide(source_layout)
        new_py_slide = PySlide(new_slide_pptx)

        for shape in source_slide_pptx.shapes:
            try:
                if shape.is_placeholder:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                        if new_py_slide.pptx_slide.shapes.title and hasattr(shape, "text"):
                            new_py_slide.pptx_slide.shapes.title.text = shape.text
                    elif hasattr(shape, "text") and shape.text and \
                         shape.placeholder_format.type in (PP_PLACEHOLDER.BODY,
                                                           PP_PLACEHOLDER.CONTENT,
                                                           PP_PLACEHOLDER.OBJECT,
                                                           PP_PLACEHOLDER.SUBTITLE,
                                                           PP_PLACEHOLDER.TEXT_BOX):
                        new_py_slide.add_text_box(shape.text,
                                                   shape.left.inches, shape.top.inches,
                                                   shape.width.inches, shape.height.inches)

                elif shape.has_text_frame and shape.text_frame.text and not shape.is_placeholder:
                    new_py_slide.add_text_box(shape.text_frame.text,
                                               shape.left.inches, shape.top.inches,
                                               shape.width.inches, shape.height.inches)

                elif hasattr(shape, 'shape_type') and shape.shape_type in (MSO_SHAPE.RECTANGLE,
                                                                            MSO_SHAPE.OVAL,
                                                                            MSO_SHAPE.LINE,
                                                                            MSO_SHAPE.ROUNDED_RECTANGLE,
                                                                            MSO_SHAPE.TRIANGLE):
                    new_added_shape = new_py_slide.add_shape(
                        shape.shape_type,
                        shape.left.inches, shape.top.inches,
                        shape.width.inches, shape.height.inches,
                        shape_name=shape.name + "_copy" if shape.name else None
                    )
                    if shape.has_text_frame and shape.text_frame.text:
                        new_added_shape.text_frame.text = shape.text_frame.text
            except Exception:
                pass

        return new_py_slide

    def save(self, filename):
        """Saves the presentation to the given filename."""
        self.presentation.save(filename)

    def set_slide_numbers_visibility(self, visible=True):
        """Attempts to set visibility of slide numbers on each slide by interacting with placeholders.

        For hiding (visible=False): Clears text from slide number placeholders.
        For showing (visible=True): This method doesn't explicitly force show if master/layout hides it,
                                   but ensures text isn't cleared if a placeholder exists.
                                   Proper enabling is best done via slide master/layout design.
        Args:
            visible (bool): True to attempt to show/ensure not hidden, False to attempt to hide.
        """
        print("INFO: Slide number visibility is best configured in the slide master/layout.")
        for i, slide_obj in enumerate(self.presentation.slides): # Renamed slide to slide_obj to avoid conflict
            slide_number_shape = None
            # Check shapes that are placeholders first
            for shape in slide_obj.placeholders: # Use slide_obj here
                if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                    slide_number_shape = shape
                    break
            # If not found in placeholders, check all shapes (less common for true slide numbers)
            if not slide_number_shape:
                for shape in slide_obj.shapes: # Use slide_obj here
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder and \
                       hasattr(shape, "placeholder_format") and \
                       shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                        slide_number_shape = shape
                        break

            if slide_number_shape:
                if not visible:
                    slide_number_shape.text_frame.clear() # Clears text, effectively hiding it.
                    print(f"INFO: Cleared text from slide number placeholder on slide {i} to hide it.")
                else:
                    # If it was cleared, this won't automatically restore it.
                    # python-pptx should fill it if the layout/master expects a slide number.
                    # We are not adding text here as it should be automatic.
                    print(f"INFO: For slide {i}, ensure layout/master enables slide numbers for placeholder to fill.")
            elif visible:
                print(f"WARNING: Slide {i} does not seem to have a slide number placeholder to make visible.")
