# presentation.py in pypptx directory

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE

# Import PySlide and constants from within the pypptx package
from .slide import PySlide
from .constants import DEFAULT_LAYOUT_REF

class PyPPT:
    def __init__(self, pptx_path=None):
        if pptx_path:
            self.presentation = Presentation(pptx_path)
        else:
            self.presentation = Presentation()
        # Store the path if provided, though it's less relevant if creating new
        self.pptx_path = pptx_path

    def add_slide(self, layout_ref=DEFAULT_LAYOUT_REF):
        """Adds a new slide to the presentation and returns its PySlide wrapper.

        Args:
            layout_ref (int or str): The index of the slide layout to use (int)
                                     or the name of the slide layout (str).
                                     Defaults to DEFAULT_LAYOUT_REF.

        Returns:
            PySlide: A wrapper for the newly added slide.

        Raises:
            ValueError: If layout_ref is a string and no layout with that name is found.
            IndexError: If layout_ref is an integer and is out of range.
            TypeError: If layout_ref is not an int or str.
        """
        slide_layout = None
        if isinstance(layout_ref, str):
            layout_name_to_find = layout_ref
            found_layout = None
            for layout in self.presentation.slide_layouts:
                if hasattr(layout, 'name') and layout.name == layout_name_to_find: # Check if layout has a name attribute
                    found_layout = layout
                    break
            if found_layout:
                slide_layout = found_layout
            else:
                available_layout_names = [l.name for l in self.presentation.slide_layouts if hasattr(l, 'name')]
                raise ValueError(
                    f"Layout with name '{layout_name_to_find}' not found. "
                    f"Available layout names are: {available_layout_names}"
                )
        elif isinstance(layout_ref, int):
            try:
                slide_layout = self.presentation.slide_layouts[layout_ref]
            except IndexError:
                raise IndexError(
                    f"Layout index {layout_ref} is out of range. "
                    f"Available layouts: {len(self.presentation.slide_layouts)}."
                )
        else:
            raise TypeError(f"layout_ref must be an integer index or a string name, not {type(layout_ref)}.")

        new_pptx_slide = self.presentation.slides.add_slide(slide_layout)
        return PySlide(new_pptx_slide)

    def get_slide(self, slide_index):
        """Gets the slide at the specified index as a PySlide instance.

        Args:
            slide_index (int): The index of the slide to retrieve.

        Returns:
            PySlide: A wrapper for the specified slide.

        Raises:
            IndexError: If slide_index is out of range.
        """
        if slide_index < 0 or slide_index >= len(self.presentation.slides):
            raise IndexError(f"Slide index {slide_index} is out of range.")
        pptx_slide = self.presentation.slides[slide_index]
        return PySlide(pptx_slide)

    @property
    def slides(self):
        """Returns a list of PySlide instances for all slides in the presentation."""
        return [PySlide(s) for s in self.presentation.slides]

    def delete_slide(self, slide_index):
        """Deletes a slide from the presentation by its index.

        Args:
            slide_index (int): The index of the slide to delete.

        Raises:
            IndexError: If slide_index is out of range.

        Note: This method manipulates internal structures of python-pptx
              and should be used with caution.
        """
        # Use len(self.presentation.slides) for public API consistency if available
        # but _sldIdLst is the direct list being manipulated.
        num_slides_in_list = len(self.presentation.slides._sldIdLst)
        if not 0 <= slide_index < num_slides_in_list:
            raise IndexError(f"Slide index {slide_index} is out of range. "
                             f"Presentation has {num_slides_in_list} slides (indices 0 to {num_slides_in_list-1}).")

        prs = self.presentation

        slide_id_entry = prs.slides._sldIdLst[slide_index]
        rId = slide_id_entry.rId

        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[slide_index]

    def move_slide(self, current_index, new_index):
        """Moves a slide from its current position to a new position.

        Args:
            current_index (int): The current index of the slide to move.
            new_index (int): The target index where the slide should be moved.
                             If new_index is greater than or equal to the number
                             of slides (after removal of the slide at current_index),
                             the slide is moved to the end.
                             If new_index is less than 0, it's treated as 0.

        Raises:
            IndexError: If current_index is out of range.

        Note: This method manipulates internal structures of python-pptx.
        """
        slides_list = self.presentation.slides._sldIdLst
        num_slides = len(slides_list)

        if not 0 <= current_index < num_slides:
            raise IndexError(f"current_index {current_index} is out of range. "
                             f"Presentation has {num_slides} slides (indices 0 to {num_slides-1}).")

        # Get the slide ID entry and remove it from its current position
        slide_id_entry_to_move = slides_list[current_index]
        del slides_list[current_index]

        # After deletion, the list is one shorter.
        # Adjust new_index to be within valid bounds for insertion into the modified list.
        num_slides_after_pop = len(slides_list) # This is num_slides - 1

        if new_index < 0:
            new_index = 0
        elif new_index > num_slides_after_pop: # If new_index is beyond the new end of list
            new_index = num_slides_after_pop   # Clamp to the end (append)

        slides_list.insert(new_index, slide_id_entry_to_move)

    def duplicate_slide(self, slide_index_to_duplicate):
        """Duplicates a slide by creating a new slide with the same layout
        and attempting to copy basic content and shapes.
        The new slide is added at the end of the presentation.

        IMPORTANT LIMITATIONS (Basic Implementation):
        - Does NOT perform a perfect, deep copy of the slide.
        - Copies text from the main title placeholder if present on both slides.
        - For other common placeholders (body, content, text box type), text content
          is copied into NEW text boxes on the duplicated slide at the same position
          and size as the original placeholder. It does not attempt to map to
          existing placeholders on the new slide by index.
        - Attempts to replicate basic auto-shapes (e.g., RECTANGLE, OVAL, LINE)
          and their text content, position, and size.
        - Does NOT copy:
            - Complex shape formatting (most fill, line, effects will be default).
            - Tables.
            - Charts.
            - Images.
            - Grouped shapes.
            - Animations, transitions, or slide master modifications.

        Args:
            slide_index_to_duplicate (int): The index of the slide to duplicate.

        Returns:
            PySlide: A PySlide wrapper for the newly created (duplicated) slide.

        Raises:
            IndexError: If slide_index_to_duplicate is out of range.
        """
        num_slides = len(self.presentation.slides)
        if not 0 <= slide_index_to_duplicate < num_slides:
            raise IndexError(f"slide_index_to_duplicate {slide_index_to_duplicate} is out of range. "
                             f"Presentation has {num_slides} slides (indices 0 to {num_slides-1}).")

        source_slide_pptx = self.presentation.slides[slide_index_to_duplicate]
        source_layout = source_slide_pptx.slide_layout

        new_slide_pptx = self.presentation.slides.add_slide(source_layout)
        new_py_slide = PySlide(new_slide_pptx) # PySlide is from .slide

        for shape in source_slide_pptx.shapes:
            try:
                if shape.is_placeholder:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE: # PP_PLACEHOLDER is from pptx.enum.shapes
                        if new_py_slide.pptx_slide.shapes.title and hasattr(shape, "text"):
                            new_py_slide.pptx_slide.shapes.title.text = shape.text
                    elif hasattr(shape, "text") and shape.text and \
                         shape.placeholder_format.type in (PP_PLACEHOLDER.BODY,
                                                           PP_PLACEHOLDER.CONTENT,
                                                           PP_PLACEHOLDER.OBJECT,
                                                           PP_PLACEHOLDER.SUBTITLE,
                                                           PP_PLACEHOLDER.TEXT_BOX):
                        # new_py_slide.add_text_box uses Inches, Pt which are in pypptx.slide
                        # MSO_SHAPE is also used by add_shape in pypptx.slide
                        new_py_slide.add_text_box(shape.text,
                                                   shape.left.inches, shape.top.inches,
                                                   shape.width.inches, shape.height.inches)

                elif shape.has_text_frame and shape.text_frame.text and not shape.is_placeholder:
                    new_py_slide.add_text_box(shape.text_frame.text,
                                               shape.left.inches, shape.top.inches,
                                               shape.width.inches, shape.height.inches)

                elif hasattr(shape, 'shape_type') and shape.shape_type in (MSO_SHAPE.RECTANGLE, # MSO_SHAPE from pptx.enum.shapes
                                                                            MSO_SHAPE.OVAL,
                                                                            MSO_SHAPE.LINE,
                                                                            MSO_SHAPE.ROUNDED_RECTANGLE,
                                                                            MSO_SHAPE.TRIANGLE):
                    new_added_shape = new_py_slide.add_shape( # new_py_slide is an instance of PySlide from .slide
                        shape.shape_type,
                        shape.left.inches, shape.top.inches,
                        shape.width.inches, shape.height.inches,
                        shape_name=shape.name + "_copy" if shape.name else None
                    )
                    if shape.has_text_frame and shape.text_frame.text:
                        new_added_shape.text_frame.text = shape.text_frame.text
            except Exception:
                # Allowing broad exception as per original code
                pass

        return new_py_slide

    def save(self, filename):
        """Saves the presentation to the given filename."""
        self.presentation.save(filename)

    def set_slide_numbers_visibility(self, visible=True):
        """Attempts to set visibility of slide numbers on each slide by interacting with placeholders.

        For hiding (visible=False): Clears text from slide number placeholders.
        For showing (visible=True): This method doesn't explicitly force show if master/layout hides it,
                                   but ensures text isn't cleared if a placeholder exists.
                                   Proper enabling is best done via slide master/layout design.
        Args:
            visible (bool): True to attempt to show/ensure not hidden, False to attempt to hide.
        """
        print("INFO: Slide number visibility is best configured in the slide master/layout.")
        for i, slide_obj in enumerate(self.presentation.slides):
            slide_number_shape = None
            # Check shapes that are placeholders first
            for shape in slide_obj.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER: # PP_PLACEHOLDER from pptx.enum.shapes
                    slide_number_shape = shape
                    break
            # If not found in placeholders, check all shapes (less common for true slide numbers)
            if not slide_number_shape:
                for shape in slide_obj.shapes:
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder and \
                       hasattr(shape, "placeholder_format") and \
                       shape.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                        slide_number_shape = shape
                        break

            if slide_number_shape:
                if not visible:
                    slide_number_shape.text_frame.clear()
                    print(f"INFO: Cleared text from slide number placeholder on slide {i} to hide it.")
                else:
                    print(f"INFO: For slide {i}, ensure layout/master enables slide numbers for placeholder to fill.")
            elif visible:
                print(f"WARNING: Slide {i} does not seem to have a slide number placeholder to make visible.")
