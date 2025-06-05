from pyppt import PyPPT
import pandas as pd
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# --- Create a new presentation ---
print("Creating a new presentation...")
new_preso = PyPPT()

# Add slides and get their PySlide instances
print("Adding slide 1 (default layout)...")
slide_one = new_preso.add_slide() # Returns PySlide
try:
    slide_one.set_title("Slide One Title (Default Layout)")
    print("Set title for slide one.")
    # Example: Add a quick bullet list to the first slide
    slide_one.add_bullet_point_box(["Item A", "Item B"], left=1, top=2, width=5, height=2)
    print("Added bullets to slide one.")
except AttributeError as e:
    print(f"Could not set title or add content for slide one: {e} (Layout might not support title or check placeholder names)")
except Exception as e:
    print(f"An error occurred with slide one: {e}")


print("\nAdding slide 2 using layout name 'Title Slide' (or 'Title' if that's the name)...")
slide_two = None # Initialize slide_two
try:
    # Common layout names. Actual names can vary by template.
    # Attempt "Title Slide" first, then "Title", then "Blank" as fallbacks for example.
    layout_name_to_try = "Title Slide"
    try:
        slide_two = new_preso.add_slide(layout_ref=layout_name_to_try)
    except ValueError:
        print(f"Layout '{layout_name_to_try}' not found, trying 'Title'...")
        layout_name_to_try = "Title"
        try:
            slide_two = new_preso.add_slide(layout_ref=layout_name_to_try)
        except ValueError:
            print(f"Layout '{layout_name_to_try}' not found, trying 'Blank'...")
            layout_name_to_try = "Blank"
            slide_two = new_preso.add_slide(layout_ref=layout_name_to_try)

    slide_two.set_title(f"Slide via Layout Name: '{layout_name_to_try}'")
    # Attempt to set subtitle, common for "Title Slide" or "Title" layouts
    try:
        slide_two.set_subtitle("Subtitle for name-based layout slide")
    except AttributeError:
        print(f"Layout '{layout_name_to_try}' does not have a subtitle placeholder.")
    print(f"Added slide using layout name '{layout_name_to_try}'.")
except Exception as e:
    print(f"Error adding slide by name: {e}. Adding by index 0 as fallback.")
    slide_two = new_preso.add_slide(layout_ref=0) # Fallback to index
    slide_two.set_title("Slide via Index 0 (Fallback)")


print("\nAdding slide 3 (layout index 1 - typically Title and Content)...")
slide_three = new_preso.add_slide(layout_ref=1) # Returns PySlide
try:
    slide_three.set_title("Slide Three (Title and Content Layout)")
    # Add a text box to this slide
    slide_three.add_text_box("Content text box on slide three.", left=1, top=2, width=5, height=2)
    print("Set title and added text box for slide three.")
except AttributeError as e:
    print(f"Could not set title or add content for slide three: {e} (Layout might not have these placeholders).")
except Exception as e:
    print(f"An error occurred with slide three: {e}")

# Save the new presentation (first version)
output_filename = "sample_presentation_initial.pptx"
print(f"\nSaving initial presentation as {output_filename}...")
new_preso.save(output_filename)
print(f"Initial presentation saved: {output_filename}")


# --- Demonstrating Further Text Manipulation on Existing Slides ---
print("\n--- Demonstrating Further Text Manipulation ---")

# Check if there are any slides to manipulate
if len(new_preso.slides) > 0:
    # Get the first slide (slide_one) using the .slides property or get_slide()
    # target_slide_ops = new_preso.slides[0]
    target_slide_ops = new_preso.get_slide(0) # Using get_slide for variety

    slide_title_for_print = "N/A"
    try:
        # Attempt to get the title text if the title shape and text exist
        if target_slide_ops.pptx_slide.shapes.title and target_slide_ops.pptx_slide.shapes.title.has_text_frame and target_slide_ops.pptx_slide.shapes.title.text_frame.text:
            slide_title_for_print = target_slide_ops.pptx_slide.shapes.title.text_frame.text
    except AttributeError: # Handles cases where title shape might not exist as expected
        pass # slide_title_for_print remains "N/A"

    print(f"\nModifying first slide (Index 0 - Current Title: '{slide_title_for_print}')...")

    print("Updating title and subtitle for first slide...")
    try:
        target_slide_ops.set_title("Updated Title for First Slide")
        target_slide_ops.set_subtitle("This is an updated subtitle for the first slide.")
        print("Title and subtitle updated for first slide.")
    except AttributeError as e:
        print(f"Could not set/update title/subtitle for first slide: {e} (Layout might not have these placeholders).")

    print("\nSetting footer text for first slide...")
    try:
        target_slide_ops.set_footer_text("Updated Footer Text for First Slide")
        print("Footer text set for first slide.")
    except AttributeError as e:
        print(f"Could not set footer for first slide: {e} (Layout might not have a footer placeholder).")

    print("\nAdding another text box to the first slide...")
    try:
        target_slide_ops.add_text_box("Another text box, added during text manipulation phase.", left=1, top=4, width=5, height=1)
        print("Another text box added to the first slide.")
    except Exception as e:
        print(f"Error adding another text box to first slide: {e}")

    print("\nAdding more bullet points to the first slide...")
    try:
        more_bullet_items = ["Additional Point 1", "Additional Point 2"]
        # This will add a new text box with bullets. If you want to add to existing, that's a different logic.
        target_slide_ops.add_bullet_point_box(more_bullet_items, left=1, top=5, width=5, height=1.5)
        print("Additional bullet point box added to the first slide.")
    except Exception as e:
        print(f"Error adding more bullet points to first slide: {e}")
else:
    print("\nSkipping text manipulation examples as no slides were added.")

# --- Demonstrating Adding DataFrame as Table ---
print("\n--- Demonstrating Adding DataFrame as Table ---")

if len(new_preso.slides) == 0:
    print("Adding a new slide for the table example...")
    table_slide = new_preso.add_slide()
    try:
        table_slide.set_title("DataFrame Table Example")
    except AttributeError:
        print("INFO: New slide for table has no title placeholder.")
else:
    print("Using the first slide for the table example.")
    table_slide = new_preso.slides[0] # Get the first slide
    try:
        if not table_slide.pptx_slide.shapes.title:
             table_slide.set_title("DataFrame Table Example (on existing slide)")
        elif not table_slide.pptx_slide.shapes.title.text:
             table_slide.set_title("DataFrame Table Example (on existing slide)")
        # If title exists and has text, we might not want to overwrite it here, or choose to.
        # For this example, we'll assume if it has text, we leave it.
    except AttributeError:
        print(f"INFO: Slide {new_preso.slides.index(table_slide)} has no title placeholder, adding table without setting/checking title here.")


# Create a sample DataFrame
data = {
    'Product Name': ['Apples', 'Bananas', 'Cherries', 'Dates'],
    'Category': ['Fruit', 'Fruit', 'Fruit', 'Fruit'],
    'Quantity': [120, 250, 75, 100],
    'Unit Price': [0.99, 0.59, 2.49, 3.00],
    'Total Value': [118.8, 147.5, 186.75, 300.00]
}
df_sample = pd.DataFrame(data)
df_sample.set_index('Product Name', inplace=True)

print("Sample DataFrame created:")
print(df_sample)

custom_headers = {
    'Category': 'Type',
    'Quantity': 'Amount (Units)',
    'Unit Price': 'Price/Unit',
    'Total Value': 'Subtotal'
}

number_formats_for_table = {
    'Unit Price': '$.2f',
    'Total Value': '$,.2f',
    'Quantity': ',d'
}

# Define some styling options
table_font_name = 'Calibri'
table_font_size = 10
num_df_columns = len(df_sample.columns)
example_col_widths = [1.5] + [1.0] * num_df_columns # Widths in Inches

header_custom_style = {
    'header_bold': True,
    'header_font_color_rgb': (255, 255, 255), # White
    'header_fill_color_rgb': (31, 73, 125)    # A dark blue
}

try:
    print("\nAdding DataFrame to slide with custom styling...")
    # Determine slide index for print message
    slide_idx_for_msg = "N/A"
    for idx, s_wrapper in enumerate(new_preso.slides):
        if s_wrapper.pptx_slide == table_slide.pptx_slide: # Compare actual pptx_slide objects
            slide_idx_for_msg = idx
            break

    table_shape = table_slide.add_table_from_dataframe(
        dataframe=df_sample,
        left=1, top=3, width=8, height=1.5, # Adjusted top position, ensure width accommodates columns
        column_labels=custom_headers,
        number_formats=number_formats_for_table,
        include_index=True,
        index_label='Product',

        # New styling options
        font_name=table_font_name,
        font_size=table_font_size,
        column_widths=example_col_widths,
        header_bold=header_custom_style['header_bold'],
        header_font_color_rgb=header_custom_style['header_font_color_rgb'],
        header_fill_color_rgb=header_custom_style['header_fill_color_rgb']
    )
    print(f"Styled DataFrame table added to slide index {slide_idx_for_msg}.")
except Exception as e:
    print(f"Error adding styled DataFrame table to slide: {e}")


# --- Demonstrating Adding and Styling Shapes ---
print("\n--- Demonstrating Adding and Styling Shapes ---")

if not new_preso.slides:
    print("Adding a new slide for shape examples...")
    # Attempt to use a "Blank" layout by name, fallback to index 6 if not found
    try:
        shape_slide = new_preso.add_slide(layout_ref="Blank")
    except ValueError:
        print("Layout 'Blank' not found, trying index 6 for a blank-like layout.")
        shape_slide = new_preso.add_slide(layout_ref=6) # Index 6 is often Blank

    try:
        shape_slide.set_title("Shape Examples")
    except AttributeError:
        print("INFO: Shape slide has no title placeholder.")
else:
    # Use the first slide for shapes
    shape_slide = new_preso.slides[0]
    print(f"Adding shapes to slide index 0...")


# Add a rectangle and name it
print("Adding a named rectangle...")
rect_shape = shape_slide.add_shape(
    MSO_SHAPE.RECTANGLE,
    left=0.5, top=4.5, width=2.5, height=1.5, # Adjusted top to avoid overlap with table
    shape_name="MyExampleRectangle"
)
print(f"Rectangle '{rect_shape.name}' added.")

# Add an oval
print("Adding an oval...")
oval_shape = shape_slide.add_shape(
    MSO_SHAPE.OVAL,
    left=3.5, top=4.5, width=2.0, height=1.0, # Adjusted top
    shape_name="MyOvalShape" # Give it a name for potential reference
)
print("Oval added.")

# Add a line
print("Adding a line shape...")
line_shape = shape_slide.add_shape(
    MSO_SHAPE.LINE, # Using MSO_SHAPE.LINE for a standard line
    left=0.5, top=6.0, width=5.0, height=0 # Adjusted top
)
print("Line shape added.")


print("\nStyling the shapes...")
# Style the rectangle by its name
print("Setting fill color for 'MyExampleRectangle'...")
try:
    shape_slide.set_shape_fill_color("MyExampleRectangle", r=173, g=216, b=230) # Light Blue
except ValueError as e:
    print(f"Could not style rectangle by name: {e}")


# Style the oval using its object reference
print("Setting line color and weight for the oval shape...")
shape_slide.set_shape_line_color(oval_shape, r=255, g=0, b=0)   # Red
shape_slide.set_shape_line_weight(oval_shape, weight_pt=3)     # 3 points thick

# Style the line by its index (assuming it's the last shape added to this slide)
# This is fragile; using the object 'line_shape' or its name (if set) is better.
# For demonstration, let's assume we know its index or give it a name.
# Let's style it by object reference for robustness:
print("Setting line color and weight for the line shape...")
shape_slide.set_shape_line_color(line_shape, r=0, g=0, b=0) # Black
shape_slide.set_shape_line_weight(line_shape, weight_pt=2)

print("Shape styling applied.")


# --- Demonstrating Adding Charts ---
print("\n--- Demonstrating Adding Charts ---")

# Get a slide to add charts to, or create a new one
if not new_preso.slides:
    print("Adding a new slide for chart examples...")
    try:
        chart_slide = new_preso.add_slide(layout_ref="Blank")
    except ValueError:
        print("Layout 'Blank' not found, using default layout for chart slide.")
        chart_slide = new_preso.add_slide() # Default layout as fallback
    try:
        chart_slide.set_title("Chart Examples")
    except AttributeError:
        print("INFO: Chart slide's layout has no title placeholder.")
else:
    # Add charts to a new slide to keep examples clean, even if other slides exist
    print("Adding a new slide for chart examples...")
    try:
        # Try to use a layout often suitable for content + title
        chart_slide = new_preso.add_slide(layout_ref="Title and Content")
    except ValueError:
        print("Layout 'Title and Content' not found, using default layout for chart slide.")
        chart_slide = new_preso.add_slide() # Default layout as fallback
    try:
        chart_slide.set_title("Chart Demonstrations")
    except AttributeError:
         print(f"INFO: Layout '{chart_slide.pptx_slide.slide_layout.name}' might not have a title placeholder.")


# Example 1: Line Chart
print("Preparing data for a Line Chart...")
line_chart_data = {
    'categories': ['Q1 Sales', 'Q2 Sales', 'Q3 Sales', 'Q4 Sales'],
    'series': [
        {'name': 'Product Alpha', 'values': [150, 200, 180, 220]},
        {'name': 'Product Beta', 'values': [120, 170, 160, 190]}
    ]
}
try:
    print("Adding a Line Chart...")
    line_chart_graphic_frame = chart_slide.add_chart(
        XL_CHART_TYPE.LINE,
        line_chart_data,
        left=0.5, top=1.5, width=4.5, height=2.5, # Inches
        chart_title="Product Sales Trends (Line)"
    )
    print("Line Chart added.")
except ValueError as e:
    print(f"Error adding Line Chart: {e}")
except Exception as e:
    print(f"An unexpected error occurred while adding Line Chart: {e}")

# Example 2: Clustered Column Chart
print("\nPreparing data for a Clustered Column Chart...")
column_chart_data = {
    'categories': ['North', 'South', 'East', 'West'],
    'series': [
        {'name': 'Y2022', 'values': [2500, 3200, 1800, 2900]},
        {'name': 'Y2023', 'values': [2800, 3000, 2000, 3100]}
    ]
}
try:
    print("Adding a Clustered Column Chart...")
    column_chart_graphic_frame = chart_slide.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        column_chart_data,
        left=5.0, top=1.5, width=4.5, height=2.5, # Adjusted left to avoid overlap
        chart_title="Regional Performance (Column)"
    )
    print("Clustered Column Chart added.")
except ValueError as e:
    print(f"Error adding Clustered Column Chart: {e}")
except Exception as e:
    print(f"An unexpected error occurred while adding Clustered Column Chart: {e}")


print("\nAttempting to set slide numbers visibility (True)...")
new_preso.set_slide_numbers_visibility(True) # This method is on PyPPT

# Re-save the presentation to include these text changes
updated_output_filename = "sample_presentation_with_text.pptx"
print(f"\nSaving presentation with all manipulations as {updated_output_filename}...")
new_preso.save(updated_output_filename)
print(f"Presentation saved: {updated_output_filename}")

print("\n--- Example Finished ---")
print(f"Please open '{output_filename}' and '{updated_output_filename}' to view the results.")

# --- Optional: Example of opening an existing presentation ---
# print("\n--- Opening an existing presentation (example) ---")
# try:
#     # To run this, first ensure 'existing_example.pptx' exists.
#     # You can create one by renaming one of the outputs from this script.
    # dummy_preso_for_opening = PyPPT()
#     # slide = dummy_preso_for_opening.add_slide()
#     # slide.set_title("Existing Presentation Example")
#     # dummy_preso_for_opening.save("existing_example.pptx")
#
#     print("Opening 'existing_example.pptx'...")
#     existing_preso = PyPPT("existing_example.pptx")
#
#     if len(existing_preso.slides) > 0:
#         print("Modifying the first slide of the existing presentation...")
#         slide_to_modify = existing_preso.slides[0] # Get PySlide for the first slide
#         slide_to_modify.set_title("Title Modified in Existing Presentation")
#         slide_to_modify.add_text_box("New text box in existing.", 1, 2, 3, 1)
#
#         print("Adding a new slide to the existing presentation...")
#         new_slide_in_existing = existing_preso.add_slide()
#         new_slide_in_existing.set_title("Newly Added Slide in Existing Presentation")
#
#         existing_output_filename = "existing_example_modified.pptx"
#         print(f"Saving modified existing presentation as {existing_output_filename}...")
#         existing_preso.save(existing_output_filename)
#         print(f"Modified existing presentation saved: {existing_output_filename}")
#     else:
#         print("'existing_example.pptx' has no slides to modify.")
#
# except FileNotFoundError:
#     print("Error: 'existing_example.pptx' not found. Please create it to run this part of the example.")
# except Exception as e:
#     print(f"Error in opening/modifying existing presentation example: {e}")
